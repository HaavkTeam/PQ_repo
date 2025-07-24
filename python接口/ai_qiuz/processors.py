import os
import tempfile
import logging
from pathlib import Path
import PyPDF2
from pptx import Presentation
import docx
import speech_recognition as sr
from pydub import AudioSegment

logger = logging.getLogger(__name__)


def pdf_to_text(file_path: str) -> str:
    """将PDF文件转换为文本"""
    logger.info(f"开始转换PDF文件: {file_path}")
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            logger.info(f"PDF文件共有 {num_pages} 页")

            for i, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                logger.debug(f"已处理第 {i + 1} 页")

        logger.info(f"PDF转换完成，提取文本长度: {len(text)}")
        return text.strip()
    except Exception as e:
        logger.error(f"PDF转换失败: {str(e)}")
        raise


def pptx_to_text(file_path: str) -> str:
    """将PPTX文件转换为文本"""
    logger.info(f"开始转换PPTX文件: {file_path}")
    text = ""
    try:
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        logger.info(f"PPTX文件共有 {num_slides} 页幻灯片")

        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            if slide_text:
                text += f"幻灯片 {i + 1}:\n{slide_text}\n"
            logger.debug(f"已处理第 {i + 1} 页幻灯片")

        logger.info(f"PPTX转换完成，提取文本长度: {len(text)}")
        return text.strip()
    except Exception as e:
        logger.error(f"PPTX转换失败: {str(e)}")
        raise


def docx_to_text(file_path: str) -> str:
    """将DOCX文件转换为文本"""
    logger.info(f"开始转换DOCX文件: {file_path}")
    text = ""
    try:
        doc = docx.Document(file_path)
        num_paras = len(doc.paragraphs)
        logger.info(f"DOCX文件共有 {num_paras} 段")

        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"

        logger.info(f"DOCX转换完成，提取文本长度: {len(text)}")
        return text.strip()
    except Exception as e:
        logger.error(f"DOCX转换失败: {str(e)}")
        raise


def txt_to_text(file_path: str) -> str:
    """读取TXT文件内容"""
    logger.info(f"开始读取TXT文件: {file_path}")
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
            logger.info(f"TXT文件读取完成，文本长度: {len(text)}")
            return text.strip()
    except Exception as e:
        logger.error(f"TXT文件读取失败: {str(e)}")
        raise


def audio_to_text(audio_path: str) -> str:
    """将音频文件转换为文本"""
    logger.info(f"开始转换音频文件: {audio_path}")
    r = sr.Recognizer()

    try:
        # 获取音频格式
        audio_format = Path(audio_path).suffix.lower()
        wav_path = audio_path

        # 处理不同音频格式
        if audio_format != '.wav':
            logger.info(f"转换音频格式: {audio_format} -> .wav")
            try:
                sound = AudioSegment.from_file(audio_path)
                wav_path = "temp_audio.wav"
                sound.export(wav_path, format="wav")
                logger.info("音频格式转换完成")
            except Exception as e:
                logger.error(f"音频格式转换失败: {str(e)}")
                raise ValueError(f"不支持的音频格式: {audio_format}")

        # 识别音频
        with sr.AudioFile(wav_path) as source:
            logger.info("开始识别音频内容...")
            audio_data = r.record(source)
            try:
                text = r.recognize_google(audio_data, language='zh-CN')
                logger.info(f"音频识别成功，文本长度: {len(text)}")
                return text.strip()
            except sr.UnknownValueError:
                logger.error("无法识别音频内容")
                raise ValueError("无法识别音频内容")
            except sr.RequestError as e:
                logger.error(f"语音识别API请求失败: {str(e)}")
                raise ConnectionError(f"语音识别服务不可用: {str(e)}")

    except Exception as e:
        logger.error(f"音频转换失败: {str(e)}")
        raise
    finally:
        # 清理临时文件
        if audio_format != '.wav' and os.path.exists(wav_path):
            os.remove(wav_path)
            logger.info("已清理临时音频文件")


def process_uploaded_file(file_content: bytes, file_type: str) -> str:
    """处理上传的文件并提取文本"""
    logger.info(f"处理上传文件，类型: {file_type}")

    # 创建临时文件
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name
        logger.info(f"创建临时文件: {tmp_path}")

    try:
        file_type = file_type.lower()

        if file_type == "pdf":
            return pdf_to_text(tmp_path)
        elif file_type in ["pptx", "ppt"]:
            return pptx_to_text(tmp_path)
        elif file_type == "docx":
            return docx_to_text(tmp_path)
        elif file_type in ["txt", "text"]:
            return txt_to_text(tmp_path)
        elif file_type in ["mp3", "wav", "flac", "aac"]:
            return audio_to_text(tmp_path)
        else:
            logger.error(f"不支持的文件格式: {file_type}")
            raise ValueError(f"不支持的文件格式: {file_type}")

    except Exception as e:
        logger.error(f"文件处理失败: {str(e)}")
        raise
    finally:
        # 确保删除临时文件
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
            logger.info(f"已删除临时文件: {tmp_path}")


