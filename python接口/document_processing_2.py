import os
import argparse
from pathlib import Path
import PyPDF2
from pptx import Presentation
import docx
import speech_recognition as sr
from pydub import AudioSegment


def pdf_to_text(file_path):
    """将PDF文件转换为文本"""
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    return text


def pptx_to_text(file_path):
    """将PPTX文件转换为文本"""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def docx_to_text(file_path):
    """将DOCX文件转换为文本"""
    doc = docx.Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def txt_to_text(file_path):
    """读取TXT文件内容"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


def audio_to_text(audio_path):
    """将音频文件转换为文本"""
    r = sr.Recognizer()

    # 处理不同音频格式
    audio_format = Path(audio_path).suffix.lower()
    if audio_format != '.wav':
        # 转换为WAV格式
        sound = AudioSegment.from_file(audio_path)
        wav_path = "temp_audio.wav"
        sound.export(wav_path, format="wav")
    else:
        wav_path = audio_path

    # 识别音频
    with sr.AudioFile(wav_path) as source:
        audio_data = r.record(source)
        try:
            text = r.recognize_google(audio_data, language='zh-CN')
        except sr.UnknownValueError:
            text = "无法识别音频内容"
        except sr.RequestError:
            text = "API请求失败"

    # 清理临时文件
    if audio_format != '.wav' and os.path.exists(wav_path):
        os.remove(wav_path)

    return text


def file_to_text(file_path):
    """根据文件类型调用相应的转换函数"""
    file_ext = Path(file_path).suffix.lower()

    if file_ext == '.pdf':
        return pdf_to_text(file_path)
    elif file_ext == '.pptx':
        return pptx_to_text(file_path)
    elif file_ext == '.docx':
        return docx_to_text(file_path)
    elif file_ext == '.txt':
        return txt_to_text(file_path)
    elif file_ext in ['.wav', '.mp3', '.flac', '.aac']:
        return audio_to_text(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {file_ext}")


def save_text(text, output_path):
    """将文本保存到文件"""
    with open(output_path, 'w', encoding='utf-8') as file:
        file.write(text)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='文件转换工具')
    parser.add_argument('input_file', type=str, help='输入文件路径')
    parser.add_argument('--output', type=str, default='output.txt', help='输出文件路径')
    args = parser.parse_args()

    try:
        # 检查文件是否存在
        if not os.path.exists(args.input_file):
            raise FileNotFoundError(f"文件不存在: {args.input_file}")

        # 转换文件
        text = file_to_text(args.input_file)

        # 保存结果
        save_text(text, args.output)
        print(f"转换成功! 结果已保存到: {args.output}")

    except Exception as e:
        print(f"转换失败: {str(e)}")